import os
import pytesseract
import docx
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.utils import clean_text, detect_file_type
from io import BytesIO
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def ingest_file(uploaded_file):
    file_extension = detect_file_type(uploaded_file.name)
    file_bytes = uploaded_file.read()
    all_text = []

    logging.info(f"Starting ingestion for file: {uploaded_file.name} (type: {file_extension})")

    if file_extension == "pdf":
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page_num, page in enumerate(doc):
                all_text.append(page.get_text())

                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes_data = base_image["image"]
                    
                    try:
                        image = Image.open(BytesIO(image_bytes_data))
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text.strip():
                            all_text.append(f"<IMAGE_TEXT>{ocr_text.strip()}</IMAGE_TEXT>")
                            logging.info(f"Extracted text from image {img_index+1} on page {page_num+1}.")
                    except Exception as e:
                        logging.warning(f"Could not process image {img_index+1} on page {page_num+1}: {e}")
            doc.close()

        except Exception as e:
            logging.error(f"Error processing PDF file {uploaded_file.name}: {e}")
            return None

    elif file_extension == "docx":
        try:
            doc = docx.Document(BytesIO(file_bytes))
            for para in doc.paragraphs:
                all_text.append(para.text)
            logging.info(f"Extracted {len(all_text)} paragraphs from DOCX.")

        except Exception as e:
            logging.error(f"Error processing DOCX file {uploaded_file.name}: {e}")
            return None

    elif file_extension == "pptx":
        try:
            prs = Presentation(BytesIO(file_bytes))
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes_data = image.blob
                            
                            pil_image = Image.open(BytesIO(image_bytes_data))
                            ocr_text = pytesseract.image_to_string(pil_image)
                            if ocr_text.strip():
                                all_text.append(f"<IMAGE_TEXT>{ocr_text.strip()}</IMAGE_TEXT>")
                                logging.info(f"Extracted text from image on slide {slide_num+1}.")
                        except Exception as e:
                            logging.warning(f"Could not process image on slide {slide_num+1}: {e}")

        except Exception as e:
            logging.error(f"Error processing PPTX file {uploaded_file.name}: {e}")
            return None

    elif file_extension in ["png", "jpg", "jpeg"]:
        try:
            image = Image.open(BytesIO(file_bytes))
            ocr_text = pytesseract.image_to_string(image)
            all_text.append(ocr_text)
            logging.info("Extracted text from standalone image file.")
        except Exception as e:
            logging.error(f"Error processing image file {uploaded_file.name}: {e}")
            return None

    else:
        logging.error(f"Unsupported file type: {file_extension}")
        raise ValueError(f"Unsupported file type: {file_extension}")

    final_text = "\n\n".join(all_text)
    return clean_text(final_text)